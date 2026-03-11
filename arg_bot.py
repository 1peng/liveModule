import numpy as np
import os
import torch
import gradio as gr
import faiss
import sys
import time
import logging
from typing import Tuple, List, Dict, Any
import product_state

# 本地向量模型依赖
from sentence_transformers import SentenceTransformer
# 云端 LLM 依赖
from openai import OpenAI

# ======================
# 日志配置
# ======================
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# ======================
# 全局变量 (混合模式状态管理)
# ======================
GLOBAL_INDEX = None
GLOBAL_METADATA = []
GLOBAL_TEXT2VEC = None
GLOBAL_LLM_CLIENT = None
embed_dim = None

# ======================
# 配置
# ======================
VECTOR_DB_PATH = "tmp/faiss_rag_db.index"
METADATA_PATH = "tmp/faiss_metadata.npy"
KNOWLEDGE_FOLDER = "knowledge"
CURRENT_PRODUCT_INDEX_KEY = "_current_product"

# --- 本地向量模型配置 ---
# 模型会自动从 HuggingFace 下载，首次运行需联网
LOCAL_EMBED_MODEL_NAME = "./checkpoints/BAAI/bge-small-zh-v1.5"
LOCAL_EMBED_MODEL_PATH = LOCAL_EMBED_MODEL_NAME

# --- 云端 LLM 配置 ---
LLM_MODEL_ONLINE = "qwen-plus"
DASHSCOPE_API_KEY = "sk-a2a7bfa3f51b41139193a3bef692d9d6"
DASHSCOPE_BASE_URL = "https://dashscope.aliyuncs.com/compatible-mode/v1"

# 检测设备 (仅用于本地向量模型)
DEVICE = "cuda" if torch.cuda.is_available() else "cpu"
print(f"🚀 向量模型运行设备：{DEVICE}")
if DEVICE == "cuda":
    print(f"📊 GPU: {torch.cuda.get_device_name(0)}")
else:
    print("⚠️ 未检测到 CUDA，向量模型将在 CPU 上运行 (速度稍慢)")


# ======================
# 1. 本地向量模型加载
# ======================
def load_local_embedding_model(model_path: str, device: str):
    """加载本地 SentenceTransformer 模型"""
    print("\n" + "=" * 50)
    print("加载本地向量模型 (Embedding)")
    print("=" * 50)

    start = time.perf_counter()
    try:
        # trust_remote_code=True 有时需要，视具体模型而定
        model = SentenceTransformer(model_path, device=device)
        end = time.perf_counter()

        dim = model.get_sentence_embedding_dimension()
        print(f"✅ 本地向量模型加载完成 ({end - start:.2f}s)")
        print(f"   模型：{model_path}")
        print(f"   维度：{dim}")
        return model, dim
    except Exception as e:
        print(f"❌ 本地向量模型加载失败：{e}")
        print("💡 提示：请确保已安装 sentence-transformers 库 (pip install sentence-transformers)")
        raise


# ======================
# 2. 云端 LLM 客户端初始化
# ======================
def load_cloud_llm_client(api_key: str, base_url: str, model_name: str):
    """初始化云端 LLM 客户端"""
    print("\n" + "=" * 50)
    print("初始化云端语言模型 (LLM)")
    print("=" * 50)

    start = time.perf_counter()
    try:
        client = OpenAI(api_key=api_key, base_url=base_url)
        end = time.perf_counter()
        print(f"✅ 云端 LLM 客户端初始化完成 ({end - start:.2f}s)")
        print(f"   模型：{model_name}")
        print(f"   地址：{base_url}")
        return client
    except Exception as e:
        print(f"❌ 云端 LLM 初始化失败：{e}")
        raise


# ======================
# 文件读取工具 (保持不变)
# ======================
def read_text_file(path: str) -> str:
    try:
        with open(path, "r", encoding="utf-8") as f:
            return f.read()
    except:
        return ""


def read_pdf(path: str) -> str:
    text = ""
    libs = [
        ("PyPDF2", lambda p: __import__("PyPDF2").PdfReader(open(p, "rb"))),
        ("pdfplumber", lambda p: __import__("pdfplumber").open(p)),
        ("fitz", lambda p: __import__("fitz").open(p))
    ]
    for name, loader in libs:
        try:
            if name == "fitz":
                doc = loader(path)
                for page in doc: text += page.get_text() + "\n"
                doc.close()
            elif name == "pdfplumber":
                with loader(path) as pdf:
                    for page in pdf.pages:
                        t = page.extract_text()
                        if t: text += t + "\n"
            else:
                reader = loader(path)
                for page in reader.pages:
                    t = page.extract_text()
                    if t: text += t + "\n"
            if text.strip(): return text
        except:
            pass
    return text


def read_docx(path: str) -> str:
    try:
        from docx import Document
        return "\n".join([p.text for p in Document(path).paragraphs])
    except:
        return ""


def read_excel(path: str) -> str:
    try:
        import pandas as pd
        return pd.read_excel(path, dtype=str).fillna("").to_string(index=False)
    except:
        return ""


def read_csv(path: str) -> str:
    try:
        import pandas as pd
        return pd.read_csv(path, dtype=str).fillna("").to_string(index=False)
    except:
        return ""


def read_pptx(path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    except:
        return ""


def load_file(path: str) -> str:
    ext = path.lower().split(".")[-1]
    if ext in ["txt", "md", "py", "json", "xml", "yaml", "yml", "html", "css", "js", "cpp", "c", "h", "java"]:
        return read_text_file(path)
    elif ext == "pdf":
        return read_pdf(path)
    elif ext == "docx":
        return read_docx(path)
    elif ext in ["xls", "xlsx"]:
        return read_excel(path)
    elif ext == "csv":
        return read_csv(path)
    elif ext in ["ppt", "pptx"]:
        return read_pptx(path)
    return ""


def load_all_docs(folder: str = KNOWLEDGE_FOLDER) -> List[Dict[str, str]]:
    if not os.path.exists(folder):
        os.makedirs(folder, exist_ok=True)
        return []
    docs = []
    for name in os.listdir(folder):
        path = os.path.join(folder, name)
        if os.path.isfile(path):
            txt = load_file(path)
            if txt and len(txt.strip()) > 10:
                docs.append({"file": name, "content": txt})
                print(f"📄 已加载：{name}")
    return docs


def load_current_product_knowledge() -> List[Dict[str, str]]:
    """加载当前商品的知识库文件"""
    _, current_product = product_state.get_current_product()
    if not current_product:
        print("⚠️ 未找到当前商品")
        return []
    
    knowledge_path = product_state.get_knowledge_path(current_product)
    if not os.path.exists(knowledge_path):
        print(f"⚠️ 知识库文件不存在: {knowledge_path}")
        return []
    
    try:
        txt = load_file(knowledge_path)
        if txt and len(txt.strip()) > 10:
            print(f"📄 已加载商品知识库: {current_product}/knowledge.txt")
            return [{"file": f"{current_product}/knowledge.txt", "content": txt}]
    except Exception as e:
        print(f"❌ 加载知识库失败: {e}")
    
    return []


def split_text(text: str, max_len: int = 350) -> List[str]:
    """将文本按句子切分成块，保持语义完整性"""
    import re
    sentences = re.split(r'([。！？\n])', text)
    sentences = [''.join(i) for i in zip(sentences[0::2], sentences[1::2] + [''])]
    sentences = [s.strip() for s in sentences if s.strip()]
    
    chunks = []
    current_chunk = ""
    
    for sentence in sentences:
        if len(current_chunk) + len(sentence) <= max_len:
            current_chunk += sentence
        else:
            if current_chunk:
                chunks.append(current_chunk)
            current_chunk = sentence
    
    if current_chunk:
        chunks.append(current_chunk)
    
    return chunks if chunks else [text[:max_len]]


# ======================
# 核心逻辑：构建索引与向量化
# ======================
def get_product_index_paths(product_name: str) -> Tuple[str, str]:
    """获取商品索引文件路径"""
    safe_name = product_name.replace('/', '_').replace('\\', '_')
    index_path = f"tmp/faiss_index_{safe_name}.index"
    metadata_path = f"tmp/faiss_metadata_{safe_name}.npy"
    return index_path, metadata_path


def build_or_load_index(text2vec_func, expected_dim: int, force_reload: bool = False):
    """构建或加载当前商品的 FAISS 索引"""
    global GLOBAL_INDEX, GLOBAL_METADATA

    _, current_product = product_state.get_current_product()
    if not current_product:
        print("⚠️ 未找到当前商品，创建空索引")
        GLOBAL_INDEX = faiss.IndexFlatL2(expected_dim)
        GLOBAL_METADATA = []
        return

    index_path, metadata_path = get_product_index_paths(current_product)
    
    if not force_reload and os.path.exists(index_path) and os.path.exists(metadata_path):
        print(f"🔍 加载已保存的向量库 ({current_product})...")
        try:
            cpu_index = faiss.read_index(index_path)
            metadata = np.load(metadata_path, allow_pickle=True).tolist()

            if cpu_index.d != expected_dim:
                print(f"⚠️ 警告：现有索引维度 ({cpu_index.d}) 与当前模型维度 ({expected_dim}) 不匹配！")
                print("   将重新构建索引...")
                raise ValueError("Dimension mismatch")

            GLOBAL_INDEX = cpu_index
            GLOBAL_METADATA = metadata
            print(f"✅ 已加载 {len(metadata)} 个向量片段")
            return
        except Exception as e:
            print(f"   加载失败或维度不匹配：{e}")

    print(f"🆕 正在构建新向量库 ({current_product})...")
    docs = load_current_product_knowledge()
    vectors = []
    metadata = []

    if not docs:
        print("⚠️ 知识库为空，创建空索引")
        GLOBAL_INDEX = faiss.IndexFlatL2(expected_dim)
        GLOBAL_METADATA = []
        return

    total_chunks = 0
    for doc in docs:
        chunks = split_text(doc["content"])
        for chunk in chunks:
            vec = text2vec_func(chunk)
            vectors.append(vec)
            metadata.append({"file": doc["file"], "text": chunk})
            total_chunks += 1
            if total_chunks % 50 == 0:
                print(f"   已处理 {total_chunks} 个片段...")

    if not vectors:
        GLOBAL_INDEX = faiss.IndexFlatL2(expected_dim)
        GLOBAL_METADATA = []
        return

    vectors_np = np.array(vectors).astype('float32')
    cpu_index = faiss.IndexFlatL2(vectors_np.shape[1])
    cpu_index.add(vectors_np)

    faiss.write_index(cpu_index, index_path)
    np.save(metadata_path, metadata)

    GLOBAL_INDEX = cpu_index
    GLOBAL_METADATA = metadata
    print(f"✅ 向量库构建完成，共 {len(vectors)} 个向量 (维度：{vectors_np.shape[1]})")


def save_current_index():
    global GLOBAL_INDEX, GLOBAL_METADATA
    if GLOBAL_INDEX is not None:
        _, current_product = product_state.get_current_product()
        if current_product:
            index_path, metadata_path = get_product_index_paths(current_product)
            faiss.write_index(GLOBAL_INDEX, index_path)
            np.save(metadata_path, GLOBAL_METADATA)
            print(f"💾 索引已保存到磁盘 ({current_product})")


def check_and_switch_product_index(text2vec_func, expected_dim: int) -> bool:
    """检查商品是否切换，如果切换则重新加载索引"""
    global GLOBAL_INDEX, GLOBAL_METADATA
    
    _, current_product = product_state.get_current_product()
    if not current_product:
        return False
    
    index_path, metadata_path = get_product_index_paths(current_product)
    
    if os.path.exists(index_path) and os.path.exists(metadata_path):
        try:
            cpu_index = faiss.read_index(index_path)
            if cpu_index.d == expected_dim:
                metadata = np.load(metadata_path, allow_pickle=True).tolist()
                GLOBAL_INDEX = cpu_index
                GLOBAL_METADATA = metadata
                print(f"✅ 切换到商品索引: {current_product}")
                return True
        except Exception as e:
            print(f"⚠️ 加载商品索引失败: {e}")
    
    print(f"🆕 为商品 {current_product} 构建新索引...")
    build_or_load_index(text2vec_func, expected_dim, force_reload=True)
    return True


# ======================
# 核心逻辑：云端 LLM 问答
# ======================
def generate_answer_cloud(context: str, question: str, client: OpenAI, model_name: str) -> str:
    prompt = f"""你是一个严格的知识库问答助手。
任务：仅根据下方的【参考资料】回答【用户问题】。

严格规则：
1. 【参考资料】中必须包含答案才能回答。
2. 如果【参考资料】中没有答案，必须且只能回复：“资料中未找到相关信息”。
3. 禁止生成问题列表。
4. 禁止猜测或编造内容。
5. 回答要直接，不要以"根据参考资料"、"根据资料"等开头。

【参考资料】：
{context}

【用户问题】：{question}

【你的回答】："""

    try:
        response = client.chat.completions.create(
            model=model_name,
            messages=[{"role": "user", "content": prompt}],
            max_tokens=600,
            temperature=0.1,
            stream=False
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        return f"⚠️ 云端模型调用失败：{str(e)}"


def rag_chat_logic(question: str, threshold: float = 0.25) -> Tuple[str, str]:
    """混合模式 RAG 问答：本地检索 + 云端生成"""
    global GLOBAL_INDEX, GLOBAL_METADATA, GLOBAL_TEXT2VEC, GLOBAL_LLM_CLIENT

    if not question.strip():
        return "请输入问题", ""

    if GLOBAL_INDEX is None or GLOBAL_INDEX.ntotal == 0:
        return "请先在【文档管理】上传文件", ""

    # 1. 本地向量化 (快速，无费用)
    q_vec = GLOBAL_TEXT2VEC(question).reshape(1, -1)

    # 2. 本地检索 (FAISS)
    k = min(5, GLOBAL_INDEX.ntotal)
    D, I = GLOBAL_INDEX.search(q_vec, k=k)

    similarities = [max(0, 1 - d * d / 2) for d in D[0]]

    # 3. 过滤与去重
    valid_indices = [(idx, sim) for idx, sim in zip(I[0], similarities) if
                     idx < len(GLOBAL_METADATA) and sim >= threshold]

    if not valid_indices:
        return "❌ 未找到高相关内容。\n建议：降低相似度阈值或上传更多文档。", ""

    valid_indices.sort(key=lambda x: x[1], reverse=True)

    top_chunks = valid_indices[:3]

    context_parts = []
    source_files = []
    for idx, sim in top_chunks:
        chunk = GLOBAL_METADATA[idx]
        context_parts.append(f"[来源：{chunk['file']}]\n{chunk['text']}")
        if chunk['file'] not in source_files:
            source_files.append(chunk['file'])

    context = "\n\n---\n\n".join(context_parts)

    # 5. 云端生成 (调用 DeepSeek)
    answer = generate_answer_cloud(context, question, GLOBAL_LLM_CLIENT, LLM_MODEL_ONLINE)

    # 6. 格式化输出
    source_text = "📚 参考来源：\n" + "\n".join([f"{i + 1}. {f}" for i, f in enumerate(source_files)])
    source_text += f"\n\n最佳匹配相似度：{top_chunks[0][1]:.2f}"

    return answer, source_text


def refresh_docs_list():
    if os.path.exists(KNOWLEDGE_FOLDER):
        files = os.listdir(KNOWLEDGE_FOLDER)
        return "\n".join([f"📄 {f}" for f in files]) if files else "📁 文件夹为空"
    return "📁 文件夹不存在"


def process_upload_logic(files):
    """处理上传：本地向量化 + 更新索引"""
    global GLOBAL_INDEX, GLOBAL_METADATA, GLOBAL_TEXT2VEC

    if not files:
        return "未选择文件"

    if not isinstance(files, list):
        files = [files]

    new_vectors = []
    new_metadata = []
    success_count = 0

    for file in files:
        try:
            f_name = file.name if hasattr(file, 'name') else "unknown"
            content = file.read() if hasattr(file, 'read') else open(f_name, 'rb').read()

            base = os.path.basename(f_name)
            save_path = os.path.join(KNOWLEDGE_FOLDER, base)

            # 处理重名
            if os.path.exists(save_path):
                name, ext = os.path.splitext(base)
                cnt = 1
                while os.path.exists(save_path):
                    save_path = os.path.join(KNOWLEDGE_FOLDER, f"{name}_{cnt}{ext}")
                    cnt += 1

            with open(save_path, "wb") as f:
                f.write(content if isinstance(content, bytes) else content.encode('utf-8'))

            txt = load_file(save_path)
            if txt and len(txt.strip()) > 10:
                chunks = split_text(txt)
                for chunk in chunks:
                    vec = GLOBAL_TEXT2VEC(chunk)  # 本地向量化
                    new_vectors.append(vec)
                    new_metadata.append({"file": os.path.basename(save_path), "text": chunk})
                success_count += 1
        except Exception as e:
            print(f"❌ 处理失败 {f_name}: {e}")

    if not new_vectors:
        return f"⚠️ 无有效内容 (成功：{success_count})"

    vecs_np = np.array(new_vectors).astype('float32')

    if GLOBAL_INDEX is None:
        GLOBAL_INDEX = faiss.IndexFlatL2(vecs_np.shape[1])
        GLOBAL_METADATA = []

    GLOBAL_INDEX.add(vecs_np)
    GLOBAL_METADATA.extend(new_metadata)
    save_current_index()

    return f"✅ 成功添加 {len(new_vectors)} 个向量片段 (来自 {success_count} 个文件)"


# ======================
# Gradio 界面
# ======================
def create_interface():
    with gr.Blocks(title="混合 RAG (本地 ARG + 云端 LLM)", theme=gr.themes.Soft()) as demo:
        gr.Markdown("# 🚀 混合架构知识库助手")
        gr.Markdown("**本地**向量化 (隐私/快速) + **云端**推理 (强大/省显存)")

        with gr.Row():
            gr.Markdown(f"🧠 **向量模型**: `BGE-Small-ZH` (本地/{DEVICE})")
            gr.Markdown(f"💬 **语言模型**: `DeepSeek-R1` (云端)")

        with gr.Tabs():
            with gr.TabItem("💬 智能问答"):
                q_input = gr.Textbox(label="提问", placeholder="输入关于文档的问题...", lines=2)
                thresh = gr.Slider(0.0, 1.0, value=0.25, step=0.05, label="相似度阈值")
                btn_ask = gr.Button("开始回答", variant="primary")

                with gr.Row():
                    out_ans = gr.Textbox(label="🤖 回答", lines=8)
                    out_src = gr.Textbox(label="📄 来源", lines=8)

                btn_ask.click(rag_chat_logic, inputs=[q_input, thresh], outputs=[out_ans, out_src])

            with gr.TabItem("📁 文档管理"):
                f_up = gr.File(label="上传文档 (PDF/TXT/DOCX等)", file_count="multiple")
                btn_up = gr.Button("📤 上传并本地向量化", variant="primary")
                status_up = gr.Textbox(label="处理状态", lines=2)
                list_docs = gr.Textbox(label="当前文档列表", value=refresh_docs_list(), lines=10)

                btn_up.click(process_upload_logic, inputs=[f_up], outputs=[status_up]).then(
                    refresh_docs_list, inputs=[], outputs=[list_docs]
                )

                gr.Button("🔄 刷新列表").click(refresh_docs_list, inputs=[], outputs=[list_docs])

            with gr.TabItem("ℹ️ 系统状态"):
                count_vec = len(GLOBAL_METADATA) if GLOBAL_METADATA else 0
                gr.Markdown(f"""
                ### 当前状态
                - **向量库大小**: {count_vec} 个片段
                - **向量模型**: 本地运行 (无 API 费用)
                - **LLM 模型**: 云端调用 (按 Token 计费)
                - **存储路径**: `{VECTOR_DB_PATH}`
                """)

    return demo


# ======================
# 主程序入口
# ======================
def main():
    global GLOBAL_TEXT2VEC, GLOBAL_LLM_CLIENT

    print("\n" + "=" * 60)
    print("启动混合 RAG 系统 (Local Embedding + Cloud LLM)")
    print("=" * 60)

    try:
        # 1. 加载本地向量模型
        embed_model, embed_dim = load_local_embedding_model(LOCAL_EMBED_MODEL_PATH, DEVICE)

        # 封装向量化函数
        def text2vec(text: str) -> np.ndarray:
            return embed_model.encode(text, convert_to_numpy=True, normalize_embeddings=True).astype('float32')

        GLOBAL_TEXT2VEC = text2vec

        # 2. 连接云端 LLM
        GLOBAL_LLM_CLIENT = load_cloud_llm_client(DASHSCOPE_API_KEY, DASHSCOPE_BASE_URL, LLM_MODEL_ONLINE)

        # 3. 构建/加载索引
        build_or_load_index(GLOBAL_TEXT2VEC, embed_dim)

        # 4. 启动 UI
        demo = create_interface()
        print("\n✅ 系统就绪！正在启动 Web 界面...")
        demo.launch(server_name="127.0.0.1", server_port=7860)

    except Exception as e:
        print(f"\n❌ 启动失败：{e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()